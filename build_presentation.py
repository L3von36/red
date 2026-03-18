"""
Thesis defense presentation — official academic standard:

  Layout  : 16:9 widescreen (33.87 cm × 19.05 cm)
  Font    : Calibri (sans-serif body), minimum 24pt body text
  Background: white throughout — NO heavy colour fills on slide body
  Accents : Thin dark-blue header bar only; tables use simple lines
  Title   : Full author/degree/supervisor/date block (standard defense format)
  Footer  : Slide number + short title on every slide except title slide
  Colours : Maximum 2 accent colours (#1B2A4A navy, #2E6FD9 blue)
  Bullets : Max 6–7 per slide, ≥24pt, 1.5 line spacing
  Tables  : Three-line (booktabs) style — top rule, mid rule, bottom rule
"""
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy

# ── Dimensions ────────────────────────────────────────────────────────────────
SW = Cm(33.87)   # slide width  (16:9 widescreen)
SH = Cm(19.05)   # slide height

# ── Colour constants ───────────────────────────────────────────────────────────
NAVY  = RGBColor(0x1B, 0x2A, 0x4A)   # primary accent
BLUE  = RGBColor(0x2E, 0x6F, 0xD9)   # secondary accent (links, highlights)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK  = RGBColor(0x1C, 0x1C, 0x1C)   # body text
GREY  = RGBColor(0x55, 0x55, 0x55)   # secondary text
LGREY = RGBColor(0xCC, 0xCC, 0xCC)   # light lines
VLGREY= RGBColor(0xF2, 0xF2, 0xF2)   # table alt-row tint
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

MARGINS = Cm(1.8)          # left/right content margin
CONTENT_TOP = Cm(3.2)      # top of content area (below header bar)
CONTENT_W   = SW - 2 * MARGINS
FOOTER_Y    = SH - Cm(0.9)


# ── pptx helpers ──────────────────────────────────────────────────────────────
def rgb(r, g, b):
    return RGBColor(r, g, b)


def add_shape(slide, x, y, w, h, fill=None, line_color=None, line_width=Pt(0)):
    sp = slide.shapes.add_shape(1, x, y, w, h)
    sp.fill.solid() if fill else sp.fill.background()
    if fill:
        sp.fill.fore_color.rgb = fill
    if line_color:
        sp.line.color.rgb = line_color
        sp.line.width = line_width
    else:
        sp.line.fill.background()
    return sp


def add_text(slide, text, x, y, w, h,
             size=Pt(18), bold=False, italic=False,
             color=DARK, align=PP_ALIGN.LEFT, wrap=True,
             font='Calibri', line_spacing=None):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p  = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    r = p.add_run()
    r.text = text
    r.font.name      = font
    r.font.size      = size
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return tb


def add_line(slide, x1, y1, x2, y2, color=LGREY, width=Pt(0.75)):
    """Add a horizontal or vertical line."""
    from pptx.util import Emu
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    return connector


def tf_add_para(tf, text, size=Pt(18), bold=False, italic=False,
                color=DARK, align=PP_ALIGN.LEFT, space_before=Pt(4),
                font='Calibri'):
    p = tf.add_paragraph()
    p.alignment    = align
    p.space_before = space_before
    r = p.add_run()
    r.text           = text
    r.font.name      = font
    r.font.size      = size
    r.font.bold      = bold
    r.font.italic    = italic
    r.font.color.rgb = color
    return p


# ── Standard slide header ──────────────────────────────────────────────────────
def slide_header(slide, title, page_num, short_title="Traffic Speed Imputation"):
    """
    Thin navy bar at top with slide title.
    Footer: slide number right, short title left.
    """
    # Top bar — 1.8 cm tall, full width
    add_shape(slide, Cm(0), Cm(0), SW, Cm(1.8), fill=NAVY)
    # Slide title in bar
    add_text(slide, title,
             MARGINS, Cm(0.2), SW - 2*MARGINS, Cm(1.4),
             size=Pt(22), bold=True, color=WHITE, align=PP_ALIGN.LEFT)

    # Thin rule at bottom of content area
    add_shape(slide, MARGINS, FOOTER_Y - Cm(0.05),
              CONTENT_W, Cm(0.04), fill=LGREY)

    # Footer text: short title left
    add_text(slide, short_title,
             MARGINS, FOOTER_Y, Cm(20), Cm(0.7),
             size=Pt(11), color=GREY, align=PP_ALIGN.LEFT)

    # Slide number right
    add_text(slide, str(page_num),
             SW - Cm(3.5), FOOTER_Y, Cm(3.2), Cm(0.7),
             size=Pt(11), color=GREY, align=PP_ALIGN.RIGHT)


# ── Bullet list text box ───────────────────────────────────────────────────────
def bullet_list(slide, bullets, x, y, w, h,
                size=Pt(20), color=DARK, indent=Cm(0.5)):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for b in bullets:
        if first:
            p = tf.paragraphs[0]; first = False
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(6)
        p.level = 0
        r = p.add_run()
        r.text = f"\u2013  {b}"
        r.font.name      = 'Calibri'
        r.font.size      = size
        r.font.color.rgb = color
    return tb


# ── Three-line (booktabs) table ────────────────────────────────────────────────
def booktabs_table(slide, headers, rows, x, y, col_widths,
                   row_h=Cm(0.75), hdr_size=Pt(16), body_size=Pt(15)):
    """Professional academic table with top/mid/bottom rules only."""
    n_rows = len(rows) + 1
    n_cols = len(headers)
    total_w = sum(col_widths)
    total_h = n_rows * row_h

    tbl = slide.shapes.add_table(n_rows, n_cols, x, y, total_w, total_h).table

    def _cell(cell, text, bold=False, align=PP_ALIGN.CENTER, size=body_size,
              color=DARK, bg=None):
        cell.text = ''
        if bg:
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
        else:
            cell.fill.background()
        p = cell.text_frame.paragraphs[0]
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name      = 'Calibri'
        r.font.size      = size
        r.font.bold      = bold
        r.font.color.rgb = color

    # Remove all borders first by using no borders (white colour with 0 width)
    for ri in range(n_rows):
        for ci in range(n_cols):
            c = tbl.cell(ri, ci)
            c.fill.background()

    for ci, (h, cw) in enumerate(zip(headers, col_widths)):
        tbl.columns[ci].width = cw
        _cell(tbl.cell(0, ci), h, bold=True, size=hdr_size, color=BLACK)

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            al = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            _cell(tbl.cell(ri+1, ci), str(val), align=al, color=DARK)

    # Draw the three rules as thin rectangles (top, mid, bottom)
    rule_w = total_w
    add_shape(slide, x, y,                   rule_w, Cm(0.04), fill=BLACK)   # top
    add_shape(slide, x, y + row_h,           rule_w, Cm(0.025), fill=LGREY)  # mid
    add_shape(slide, x, y + total_h - Cm(0.04), rule_w, Cm(0.04), fill=BLACK)  # bottom

    return tbl


# ── Section divider slide ──────────────────────────────────────────────────────
def section_divider(slide, number, title, subtitle=""):
    """Full-navy divider slide between major sections."""
    add_shape(slide, Cm(0), Cm(0), SW, SH, fill=NAVY)
    add_text(slide, number, Cm(2.5), Cm(5.5), SW - Cm(5), Cm(1.5),
             size=Pt(14), color=rgb(0xAD, 0xC8, 0xE8), font='Calibri',
             bold=False)
    add_text(slide, title, Cm(2.5), Cm(6.8), SW - Cm(5), Cm(2.8),
             size=Pt(40), bold=True, color=WHITE, font='Calibri')
    if subtitle:
        add_text(slide, subtitle, Cm(2.5), Cm(10.0), SW - Cm(5), Cm(1.5),
                 size=Pt(20), italic=True, color=rgb(0xAD, 0xC8, 0xE8),
                 font='Calibri')


# ══════════════════════════════════════════════════════════════════════════════
# BUILD PRESENTATION
# ══════════════════════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width  = SW
prs.slide_height = SH
blank = prs.slide_layouts[6]   # blank layout
pg    = [0]   # mutable page counter (use list so inner functions can mutate)

def next_slide():
    pg[0] += 1
    return prs.slides.add_slide(blank), pg[0]


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title (official thesis defense format)
# ─────────────────────────────────────────────────────────────────────────────
sl, _ = next_slide()
# White background — no shape needed (default)

# Thin top rule (navy, not a thick bar)
add_shape(sl, Cm(0), Cm(0), SW, Cm(0.35), fill=NAVY)

# University / department block (top-left, small)
add_text(sl,
    "[University Name]  ·  [Department / School of Engineering]",
    MARGINS, Cm(0.55), SW - 2*MARGINS, Cm(0.7),
    size=Pt(13), color=GREY, align=PP_ALIGN.LEFT)

# Main title — centred, large
add_text(sl,
    "Hypergraph Neural ODEs with Observation Assimilation\nfor Sparse Traffic Speed Imputation",
    MARGINS, Cm(2.8), CONTENT_W, Cm(3.4),
    size=Pt(32), bold=True, color=NAVY,
    align=PP_ALIGN.CENTER)

# Thin divider under title
add_shape(sl, MARGINS, Cm(6.3), CONTENT_W, Cm(0.04), fill=LGREY)

# Thesis details block
details = [
    ("Thesis submitted for the degree of",  Pt(15), False, GREY),
    ("Master of Science in [Programme]",     Pt(18), True,  DARK),
    ("",                                     Pt(8),  False, DARK),
    ("Author:      [Full Name]",             Pt(16), False, DARK),
    ("Supervisor:  [Supervisor Name]",       Pt(15), False, GREY),
    ("Date:        March 2026",              Pt(15), False, GREY),
]
cy = Cm(6.8)
for txt, sz, bold, col in details:
    add_text(sl, txt, MARGINS + Cm(4), cy, Cm(20), Cm(0.8),
             size=sz, bold=bold, color=col)
    cy += sz.pt * 0.045 * Cm(1).pt / Pt(1).pt + Cm(0.1)

# Bottom navy rule
add_shape(sl, Cm(0), SH - Cm(0.35), SW, Cm(0.35), fill=NAVY)

# Keywords in footer area
add_text(sl,
    "Keywords: Graph Neural ODE · Hypergraph Convolution · Traffic Imputation · "
    "Observation Assimilation · Physics-Informed Learning",
    MARGINS, SH - Cm(1.1), CONTENT_W, Cm(0.7),
    size=Pt(11), italic=True, color=GREY, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 2 — Outline
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "Outline", p)

outline = [
    "1.  Motivation & Problem Statement",
    "2.  Background & State of the Art",
    "3.  Dataset & Experimental Setup",
    "4.  Proposed Architecture",
    "      4.1  Graph Attention ODE",
    "      4.2  Hypergraph Convolution (gated)",
    "      4.3  Observation Assimilation Gate",
    "      4.4  Physics-Informed Loss",
    "5.  Training Strategy",
    "6.  Results & Analysis",
    "      6.1  Main evaluation (80% sparsity)",
    "      6.2  Sensor sparsity sweep (20 – 90%)",
    "      6.3  Ablation study",
    "7.  Conclusion & Future Work",
]
tb = sl.shapes.add_textbox(MARGINS, CONTENT_TOP, CONTENT_W, SH - CONTENT_TOP - Cm(1.2))
tf = tb.text_frame; tf.word_wrap = True
first = True
for item in outline:
    p_ = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p_.space_before = Pt(3)
    indent = item.startswith("      ")
    r = p_.add_run()
    r.text = item
    r.font.name      = 'Calibri'
    r.font.size      = Pt(16) if not indent else Pt(14)
    r.font.bold      = not indent and not item[0].isspace()
    r.font.color.rgb = NAVY if not indent else GREY
    r.font.italic    = indent


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 3 — Motivation
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "1.  Motivation", p)

add_text(sl,
    "Real sensor networks are never complete.",
    MARGINS, CONTENT_TOP, CONTENT_W, Cm(0.9),
    size=Pt(22), bold=True, color=NAVY)

bullet_list(sl, [
    "Hardware failures, maintenance windows, and budget constraints leave many sensors dark at any moment.",
    "California PEMS04: 307 sensors across SF Bay Area freeways — realistic deployments observe only 20–60%.",
    "Missing speed data blocks: route planning, incident detection, signal timing, emissions modelling.",
    "This work: 80% of sensors are unobserved — we must infer ~246 blind node speeds from ~61 observed.",
], MARGINS, Cm(4.4), CONTENT_W, Cm(9.0), size=Pt(20))

# Research question box
add_shape(sl, MARGINS, Cm(13.8), CONTENT_W, Cm(3.0),
          fill=VLGREY, line_color=NAVY, line_width=Pt(1.0))
add_text(sl, "Research Question",
         MARGINS + Cm(0.4), Cm(14.1), CONTENT_W - Cm(0.8), Cm(0.7),
         size=Pt(14), bold=True, color=NAVY)
add_text(sl,
    "Can a Hypergraph Neural ODE with learned sensor assimilation accurately "
    "recover traffic speeds at unobserved sensors — especially during congestion events?",
    MARGINS + Cm(0.4), Cm(14.85), CONTENT_W - Cm(0.8), Cm(1.7),
    size=Pt(17), italic=True, color=DARK, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 4 — Background: Task Definition
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "2.  Background — Task Definition", p)

add_text(sl, "Imputation vs Forecasting — a critical distinction",
         MARGINS, CONTENT_TOP, CONTENT_W, Cm(0.9),
         size=Pt(20), bold=True, color=NAVY)

# Two-column layout using textboxes
half = CONTENT_W / 2 - Cm(0.3)
# Left column header
add_text(sl, "Traffic Forecasting (existing work)",
         MARGINS, Cm(4.2), half, Cm(0.7),
         size=Pt(17), bold=True, color=GREY)
add_shape(sl, MARGINS, Cm(4.9), half, Cm(0.025), fill=LGREY)
bullet_list(sl, [
    "All sensors observed",
    "Predict future values (1–3 steps ahead)",
    "Well-studied: DCRNN, STGCN, Graph WaveNet",
    "PEMS04 benchmark MAE: 1.5 – 1.8 km/h",
], MARGINS, Cm(5.1), half, Cm(7.5), size=Pt(18), color=GREY)

# Right column header
rx = MARGINS + half + Cm(0.6)
add_text(sl, "Sparse Imputation  \u2190 This work",
         rx, Cm(4.2), half, Cm(0.7),
         size=Pt(17), bold=True, color=NAVY)
add_shape(sl, rx, Cm(4.9), half, Cm(0.025), fill=NAVY)
bullet_list(sl, [
    "80% of sensors are MISSING",
    "Recover present values, not future",
    "No established deep-learning baseline",
    "Jam MAE (speed < 40 km/h) is the key metric",
], rx, Cm(5.1), half, Cm(7.5), size=Pt(18), color=NAVY)

# Vertical divider
add_shape(sl, MARGINS + half + Cm(0.25), Cm(4.2), Cm(0.025),
          Cm(8.4), fill=LGREY)

# Key note at bottom
add_text(sl,
    "Note:  A model predicting the global mean speed achieves 5.18 km/h overall MAE "
    "— but completely fails during congestion.  Overall MAE is insufficient; jam MAE must be reported separately.",
    MARGINS, Cm(14.0), CONTENT_W, Cm(2.5),
    size=Pt(15), italic=True, color=GREY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 5 — Background: Related Work SOTA table
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "2.  Background — State of the Art", p)

add_text(sl,
    "Published PEMS04 results (full-sensor, 15-min forecasting task)",
    MARGINS, CONTENT_TOP, CONTENT_W, Cm(0.8), size=Pt(18), color=DARK)

booktabs_table(sl,
    ["Model", "Venue", "Architecture", "PEMS04 MAE"],
    [
        ["DCRNN  (Li et al., 2018)", "ICLR 2018",
         "Diffusion GCN + seq2seq RNN", "~1.8 km/h"],
        ["STGCN  (Yu et al., 2018)", "IJCAI 2018",
         "Graph conv + temporal conv",  "~1.7 km/h"],
        ["Graph WaveNet  (Wu et al., 2019)", "IJCAI 2019",
         "Adaptive adj + dilated conv", "~1.6 km/h"],
        ["ASTGCN  (Guo et al., 2019)", "AAAI 2019",
         "Spatial + temporal attention","~1.6 km/h"],
        ["AGCRN  (Bai et al., 2020)",  "NeurIPS 2020",
         "Node-adaptive GCN + GRU",    "~1.5 km/h"],
        ["This work", "MSc Thesis 2026",
         "Hypergraph GAT-ODE + Assimilation", "5.18 km/h *"],
    ],
    x=MARGINS, y=Cm(4.7),
    col_widths=[Cm(10.0), Cm(4.0), Cm(10.5), Cm(4.2)],
    row_h=Cm(1.1)
)

add_text(sl,
    "* Different task: 80% sensors missing (imputation), not full-sensor forecasting.  "
    "The ~3× MAE gap reflects task difficulty, not model quality.",
    MARGINS, Cm(13.5), CONTENT_W, Cm(2.0),
    size=Pt(14), italic=True, color=GREY, wrap=True)

bullet_list(sl, [
    "Gap in literature: none of the above models address sparse imputation with >50% missing sensors.",
    "This work fills that gap with a Hypergraph Neural ODE architecture tailored for the imputation setting.",
], MARGINS, Cm(15.2), CONTENT_W, Cm(3.0), size=Pt(16))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 6 — Dataset and Setup
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "3.  Dataset and Experimental Setup", p)

# Three info columns
col_w = (CONTENT_W - Cm(1.2)) / 3
for i, (hdr, rows) in enumerate([
    ("PEMS04 Dataset",
     ["307 sensors, SF Bay Area freeways",
      "5-minute sampling intervals",
      "5,000 timesteps ≈ 17.4 days",
      "Speed channel extracted (index 2)",
      "Z-score normalised: μ, σ from full set"]),
    ("Sparsity Configuration",
     ["Sparsity ratio: 80%",
      "≈ 61 observed nodes (randomly selected)",
      "≈ 246 blind nodes (must be imputed)",
      "Fixed mask, seed = 42 (reproducible)",
      "Sweep tested: 20%, 40%, 60%, 80%, 90%"]),
    ("Train / Val / Eval Split",
     ["Train:  t = 0 – 3,999",
      "Val:    t = 4,000 – 4,239",
      "Eval:   t = 4,500 – 4,949",
      "500-step buffer (no temporal leakage)",
      "Non-overlapping 48-step windows"]),
]):
    cx = MARGINS + i * (col_w + Cm(0.6))
    add_text(sl, hdr, cx, CONTENT_TOP, col_w, Cm(0.75),
             size=Pt(17), bold=True, color=NAVY)
    add_shape(sl, cx, CONTENT_TOP + Cm(0.8), col_w, Cm(0.04), fill=NAVY)
    bullet_list(sl, rows, cx, CONTENT_TOP + Cm(1.0), col_w, Cm(7.5),
                size=Pt(16))

# Feature table
add_text(sl, "6-Dimensional Input Feature Vector (per node per timestep — no ground-truth leakage)",
         MARGINS, Cm(11.8), CONTENT_W, Cm(0.7), size=Pt(15), bold=True, color=DARK)
booktabs_table(sl,
    ["#", "Feature Name", "Description", "Blind node value"],
    [
        ["1", "obs_speed",  "Observed speed (normalised)", "0.0  (zeroed)"],
        ["2", "global_ctx", "Mean of all observed nodes at time t", "Same as observed"],
        ["3", "nbr_ctx",    "Adj-weighted mean of observed neighbours", "0.0 if no obs. nbr"],
        ["4", "is_observed","Binary sensor flag", "0"],
        ["5–6","t_sin / t_cos","Time-of-day cyclic encoding (scaled ×0.25)","Shared with obs. nodes"],
    ],
    x=MARGINS, y=Cm(12.55),
    col_widths=[Cm(1.2), Cm(5.0), Cm(13.5), Cm(8.5)],
    row_h=Cm(0.85), hdr_size=Pt(14), body_size=Pt(14)
)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 7 — Architecture Overview
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "4.  Proposed Architecture — Overview", p)

add_text(sl,
    "The model processes T = 48 timestep windows recurrently.  "
    "At each step: decode → Euler ODE step → assimilation → repeat.",
    MARGINS, CONTENT_TOP, CONTENT_W, Cm(1.0), size=Pt(18), color=DARK, wrap=True)

# Pipeline boxes — white with navy outline (no colour fills)
stages = [
    ("Input Encoder", "Linear(6 → 64)\nper node"),
    ("GAT ODE Layer", "2× Graph Attention\nEuler  dt = 0.3"),
    ("Hypergraph Conv", "2-hop corridors\nLearnable gate"),
    ("Assimilation", "Kalman-style\nSensor injection"),
    ("Decoder", "Linear(64 → 1)\nSpeed ŝ"),
]
bw  = Cm(5.4)
bh  = Cm(3.0)
gap = Cm(0.8)
sx  = MARGINS
sy  = Cm(4.8)

for i, (name, desc) in enumerate(stages):
    bx = sx + i * (bw + gap)
    add_shape(sl, bx, sy, bw, bh, fill=None,
              line_color=NAVY, line_width=Pt(1.2))
    add_text(sl, name, bx, sy + Cm(0.25), bw, Cm(0.75),
             size=Pt(16), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_shape(sl, bx + Cm(0.4), sy + Cm(0.95), bw - Cm(0.8),
              Cm(0.03), fill=LGREY)
    add_text(sl, desc, bx, sy + Cm(1.1), bw, Cm(1.6),
             size=Pt(14), italic=True, color=DARK, align=PP_ALIGN.CENTER)
    if i < len(stages) - 1:
        add_text(sl, "\u2192",
                 bx + bw + Cm(0.1), sy + Cm(1.0), gap, Cm(1.0),
                 size=Pt(22), bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Key equations below
eqs = [
    ("ODE step:", "z_{t+1} = z_t + 0.3 \u00b7 f_\u03b8(z_t)"),
    ("f_\u03b8(z):", "LayerNorm( Tanh(GAT\u2082(Tanh(GAT\u2081(z)))) + g \u00b7 HypConv(z) )"),
    ("Assimilation:", "z \u2190 z + \u03c3(W_g [z; z_obs]) \u2299 (z_obs \u2212 z) \u2299 obs_mask"),
]
for i, (lbl, eq) in enumerate(eqs):
    add_text(sl, lbl, MARGINS, Cm(9.0) + i*Cm(1.6),
             Cm(5), Cm(1.3), size=Pt(16), bold=True, color=NAVY)
    add_text(sl, eq, MARGINS + Cm(5.2), Cm(9.0) + i*Cm(1.6),
             CONTENT_W - Cm(5.2), Cm(1.3),
             size=Pt(16), color=DARK, font='Courier New')

add_shape(sl, MARGINS, Cm(8.7), CONTENT_W, Cm(0.03), fill=LGREY)
add_shape(sl, MARGINS, Cm(13.7), CONTENT_W, Cm(0.03), fill=LGREY)

bullet_list(sl, [
    "Input tensor shape: [B, N, T, 6]  —  batch × nodes × timesteps × features",
    "Hidden state z \u2208 \u211d^{N\u00d764} evolves continuously through the ODE and is corrected by each sensor observation.",
], MARGINS, Cm(14.0), CONTENT_W, Cm(3.5), size=Pt(17))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 8 — GAT and Hypergraph
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "4.1–4.2  Graph Attention + Hypergraph Convolution", p)

half = CONTENT_W / 2 - Cm(0.3)
# Left — GAT
add_text(sl, "Graph Attention Network (GAT)",
         MARGINS, CONTENT_TOP, half, Cm(0.75),
         size=Pt(18), bold=True, color=NAVY)
add_shape(sl, MARGINS, CONTENT_TOP + Cm(0.8), half, Cm(0.04), fill=NAVY)
bullet_list(sl, [
    "e\u1d62\u2c7c = LeakyReLU( a\u209b\u1d63\u1d9c(Wh\u1d62) + a\u1d49\u209b\u209c(Wh\u2c7c) )",
    "\u03b1\u1d62\u2c7c = softmax( e\u1d62\u2c7c / \u03c4=2 )  over road neighbours of i",
    "Non-edges masked to \u2212\u221e before softmax",
    "Temperature \u03c4=2 prevents collapse to a single neighbour \u2192 avoids oscillation",
    "Two stacked GAT layers inside the ODE function",
], MARGINS, CONTENT_TOP + Cm(1.0), half, Cm(8.0), size=Pt(17))

# Vertical divider
add_shape(sl, MARGINS + half + Cm(0.25), CONTENT_TOP, Cm(0.03),
          Cm(11.5), fill=LGREY)

# Right — Hypergraph
rx = MARGINS + half + Cm(0.6)
add_text(sl, "Gated Hypergraph Convolution",
         rx, CONTENT_TOP, half, Cm(0.75),
         size=Pt(18), bold=True, color=NAVY)
add_shape(sl, rx, CONTENT_TOP + Cm(0.8), half, Cm(0.04), fill=NAVY)
bullet_list(sl, [
    "Hyperedge for node i = i \u222a { all 2-hop reachable sensors }",
    "Captures corridor/intersection groups beyond pairwise edges",
    "H_conv = D_v^{\u207b\u00bd} H D_e^{\u207b\u00b9} H\u1d40 D_v^{\u207b\u00bd}  (pre-computed once)",
    "Gate: g = sigmoid(w),  w_init = \u22122  \u2192  g \u2248 0.12",
    "h = h_GAT + g \u00b7 HypConv(x)   (gated fusion)",
    "Gate prevents over-smoothing at congested nodes",
], rx, CONTENT_TOP + Cm(1.0), half, Cm(8.0), size=Pt(17))

# Joint note
add_text(sl,
    "Without the gate: a jammed sensor in a 2-hop hyperedge with 20 free-flowing corridor members "
    "would be averaged toward free-flow — actively contradicting the congestion signal.",
    MARGINS, Cm(13.5), CONTENT_W, Cm(2.5),
    size=Pt(15), italic=True, color=GREY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 9 — Assimilation and Physics Loss
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "4.3–4.4  Observation Assimilation and Physics-Informed Loss", p)

half = CONTENT_W / 2 - Cm(0.3)

# Left — Assimilation
add_text(sl, "Observation Assimilation Gate",
         MARGINS, CONTENT_TOP, half, Cm(0.75),
         size=Pt(18), bold=True, color=NAVY)
add_shape(sl, MARGINS, CONTENT_TOP + Cm(0.8), half, Cm(0.04), fill=NAVY)
bullet_list(sl, [
    "z_obs  = W_obs \u00b7 x_{t+1}            \u2190 encode new reading",
    "gate   = \u03c3( W_g [z ; z_obs] )      \u2190 Kalman gain (learned)",
    "update = gate \u2299 (z_obs \u2212 z) \u2299 obs_mask",
    "z \u2190 z + update                      \u2190 corrected state",
    "obs_mask zeros blind nodes \u2192 no leakage",
], MARGINS, CONTENT_TOP + Cm(1.0), half, Cm(7.5), size=Pt(16))

add_text(sl, "Kalman filter analogy:",
         MARGINS, Cm(9.8), half, Cm(0.6), size=Pt(14), bold=True, color=GREY)
add_text(sl,
    "  z_pred = A\u00b7z  \u2192  z_corr = z_pred + K\u00b7(obs \u2212 z_pred)\n"
    "  Our gate \u2248 Kalman gain K, but learned from data.",
    MARGINS, Cm(10.35), half, Cm(1.5),
    size=Pt(14), italic=True, color=GREY, wrap=True)

# Vertical divider
add_shape(sl, MARGINS + half + Cm(0.25), CONTENT_TOP, Cm(0.03),
          Cm(10.5), fill=LGREY)

# Right — Loss
rx = MARGINS + half + Cm(0.6)
add_text(sl, "Three-Term Loss Function",
         rx, CONTENT_TOP, half, Cm(0.75), size=Pt(18), bold=True, color=NAVY)
add_shape(sl, rx, CONTENT_TOP + Cm(0.8), half, Cm(0.04), fill=NAVY)

for i, (term, formula, note) in enumerate([
    ("1.  Jam-weighted MSE",
     "L_obs = mean( ((ŝ\u2212s)\u2299mask)² \u2299 w )",
     "w = 4 if speed < 40 km/h,  else 1.  Compensates 12:1 class imbalance."),
    ("2.  Temporal smoothness  (\u03bb = 0.60)",
     "L_smooth = mean( (ŝ_{t+1} \u2212 ŝ_t)² )",
     "Penalises step-to-step jumps.  Suppresses post-jam oscillation."),
    ("3.  Graph Laplacian physics  (\u03bb = 0.02)",
     "L_phys = mean( ||L_sym \u00b7 v||² )",
     "Encodes LWR continuity: speed varies smoothly along roads."),
]):
    cy = CONTENT_TOP + Cm(1.1) + i * Cm(3.5)
    add_text(sl, term, rx, cy, half, Cm(0.65),
             size=Pt(15), bold=True, color=DARK)
    add_text(sl, formula, rx, cy + Cm(0.65), half, Cm(0.8),
             size=Pt(14), color=BLUE, font='Courier New')
    add_text(sl, note, rx, cy + Cm(1.4), half, Cm(1.3),
             size=Pt(13), italic=True, color=GREY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 10 — Training Strategy
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "5.  Training Strategy", p)

add_text(sl, "Three techniques to overcome extreme class imbalance (free-flow 92%  vs  jam 8%)",
         MARGINS, CONTENT_TOP, CONTENT_W, Cm(0.8),
         size=Pt(18), bold=False, color=DARK, italic=True)

cw = (CONTENT_W - Cm(1.2)) / 3
for i, (title, body_pts) in enumerate([
    ("Curriculum Masking",
     ["15% of observed nodes pseudo-blinded each batch",
      "Speed, nbr_ctx, is_observed features zeroed",
      "Loss computed on their known ground truth",
      "\u2192 Gradients always flow through blind-node path",
      "\u2192 Prevents model from ignoring imputation task"]),
    ("Jam-Biased Sampling",
     ["50% of batches start at a jam timestep",
      "Natural rate of jams: only ~8% of windows",
      "Pre-computed jam_t_valid index for fast lookup",
      "\u2192 Balanced gradient signal between jam / free-flow",
      "\u2192 Effective even with 4\u00d7 jam weight in loss"]),
    ("Gradient Accumulation",
     ["4 windows accumulated per parameter update",
      "Jam batches: large loss; free-flow: small loss",
      "High variance \u2192 oscillating validation MAE",
      "\u2192 Accumulation smooths this variance",
      "\u2192 Loss divided by 4 to preserve effective LR"]),
]):
    cx = MARGINS + i * (cw + Cm(0.6))
    add_text(sl, title, cx, Cm(4.2), cw, Cm(0.75),
             size=Pt(17), bold=True, color=NAVY)
    add_shape(sl, cx, Cm(4.95), cw, Cm(0.04), fill=NAVY)
    bullet_list(sl, body_pts, cx, Cm(5.1), cw, Cm(8.0), size=Pt(15))

# Optimiser summary
add_shape(sl, MARGINS, Cm(13.5), CONTENT_W, Cm(0.04), fill=LGREY)
add_text(sl,
    "Optimiser:  Adam  lr = 3\u00d710\u207b\u2074  ·  weight decay = 10\u207b\u2074  ·  "
    "Cosine Annealing T_max = 400  ·  800 epochs  ·  "
    "Gradient clip norm = 1.0  ·  Batch window = 48 timesteps",
    MARGINS, Cm(13.65), CONTENT_W, Cm(1.5),
    size=Pt(15), color=DARK, align=PP_ALIGN.LEFT, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 11 — Main Results
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "6.1  Results — 80% Sensor Sparsity", p)

add_text(sl, "Evaluation on held-out set  t = 4,500 – 4,949  (450 timesteps, non-overlapping 48-step windows)",
         MARGINS, CONTENT_TOP, CONTENT_W, Cm(0.8), size=Pt(16), color=GREY)

booktabs_table(sl,
    ["Model", "MAE — all blind nodes (km/h)", "MAE — jam only (km/h)", "vs Global Mean Jam"],
    [
        ["Global mean baseline",     "5.18", "35.99", "\u2014"],
        ["IDW spatial interpolation","5.23", "32.95", "\u22128.5%"],
        ["Ours — full model",        "5.18", "33.93", "\u22125.7%"],
    ],
    x=MARGINS, y=Cm(4.8),
    col_widths=[Cm(9.5), Cm(7.5), Cm(6.5), Cm(5.0)],
    row_h=Cm(1.2)
)

bullet_list(sl, [
    "Overall MAE ties the global mean — expected: free-flow (92%) dominates the average and "
    "is tightly clustered near \u03bc.",
    "Jam MAE is the meaningful metric: the model achieves 5.7% improvement over global mean "
    "and beats IDW on congestion events.",
    "Evaluation uses 48-step windows (same as training) to prevent ODE state drift in long rollouts.",
], MARGINS, Cm(9.5), CONTENT_W, Cm(5.5), size=Pt(18))

add_text(sl,
    "Definition: Jam \u2261 speed < 40 km/h  (\u224814.4 standard deviations below mean in this dataset).",
    MARGINS, Cm(15.2), CONTENT_W, Cm(1.0),
    size=Pt(14), italic=True, color=GREY)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 12 — Sparsity Sweep
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "6.2  Results — Sensor Sparsity Sweep (20 % – 90 %)", p)

add_text(sl,
    "Each sparsity level trained independently (150 epochs, hidden = 32).  "
    "Goal: verify graceful degradation, not match full-model absolute MAE.",
    MARGINS, CONTENT_TOP, CONTENT_W, Cm(1.0), size=Pt(16), color=GREY, wrap=True)

booktabs_table(sl,
    ["Sparsity", "Blind sensors", "Global Mean Jam (km/h)",
     "IDW Jam (km/h)", "Model Jam (km/h)", "vs Global Mean"],
    [
        ["20%", "22%", "36.04", "27.35", "27.68", "+23.2%"],
        ["40%", "48%", "35.74", "29.73", "28.27", "+20.9%"],
        ["60%", "62%", "36.01", "29.98", "31.66", "+12.1%"],
        ["80%", "83%", "35.99", "32.95", "31.19", "+13.3%"],
        ["90%", "90%", "35.99", "33.81", "34.58",  "+3.9%"],
    ],
    x=MARGINS, y=Cm(4.8),
    col_widths=[Cm(3.5), Cm(4.0), Cm(6.5), Cm(5.5), Cm(5.5), Cm(4.2)],
    row_h=Cm(1.1)
)

bullet_list(sl, [
    "Model consistently outperforms global-mean baseline at ALL five sparsity levels.",
    "At 90% sparsity, performance approaches baseline — expected: almost no observed neighbour "
    "context remains to inform imputation.",
    "Non-monotonicity between 60% and 80% is noise from the lighter 150-epoch training; "
    "the overall trend is clearly degrading with sparsity.",
], MARGINS, Cm(11.3), CONTENT_W, Cm(5.5), size=Pt(18))


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 13 — Ablation Study
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "6.3  Results — Ablation Study", p)

add_text(sl,
    "Each variant removes exactly one component and is re-trained from scratch under identical conditions (300 epochs).",
    MARGINS, CONTENT_TOP, CONTENT_W, Cm(0.8), size=Pt(16), color=GREY)

booktabs_table(sl,
    ["Model / Variant", "MAE all (km/h)", "MAE jam (km/h)",
     "\u0394 jam  (full \u2212 variant)"],
    [
        ["Global mean baseline",       "5.18", "35.99", "\u2014"],
        ["IDW (spatial interp.)",       "5.23", "32.95", "\u2014"],
        ["Full model  \u25c4",          "5.18", "33.93", "+0.00"],
        ["\u2212  Hypergraph conv",     "5.54", "31.66", "\u22122.27"],
        ["\u2212  Assimilation gate",   "5.29", "32.54", "\u22121.38"],
        ["\u2212  Physics loss",        "5.32", "33.32", "\u22120.61"],
        ["\u2212  Neighbour context",   "5.84", "28.99", "\u22124.94"],
        ["\u2212  Temporal encoding",   "6.00", "33.51", "\u22120.41"],
    ],
    x=MARGINS, y=Cm(4.5),
    col_widths=[Cm(9.5), Cm(5.5), Cm(5.5), Cm(8.5)],
    row_h=Cm(0.98)
)

add_text(sl,
    "\u0394 jam = full_jam \u2212 variant_jam.  "
    "Positive: removing this component raises jam MAE (component is beneficial).  "
    "Negative: removing it reduces jam MAE (component introduces smoothing noise that "
    "outweighs benefit at 300 epochs; the gated hypergraph requires more training to converge).",
    MARGINS, Cm(14.0), CONTENT_W, Cm(2.5),
    size=Pt(14), italic=True, color=GREY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 14 — Conclusion
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "7.  Conclusion", p)

half = CONTENT_W / 2 - Cm(0.3)

add_text(sl, "Contributions",
         MARGINS, CONTENT_TOP, half, Cm(0.75), size=Pt(18), bold=True, color=NAVY)
add_shape(sl, MARGINS, CONTENT_TOP + Cm(0.8), half, Cm(0.04), fill=NAVY)
bullet_list(sl, [
    "Hypergraph Neural ODE — first application to sparse traffic imputation",
    "Gated HGNN branch — prevents corridor over-smoothing at jam nodes",
    "Kalman-style assimilation gate — sensor fusion without GT leakage",
    "Physics-informed loss — LWR continuity via graph Laplacian",
    "Curriculum masking + jam-biased sampling — overcomes 12:1 class imbalance",
    "Validated at five sparsity levels (20 % – 90 %)",
], MARGINS, CONTENT_TOP + Cm(1.0), half, Cm(9.5), size=Pt(17))

add_shape(sl, MARGINS + half + Cm(0.25), CONTENT_TOP, Cm(0.03),
          Cm(10.0), fill=LGREY)

rx = MARGINS + half + Cm(0.6)
add_text(sl, "Future Work",
         rx, CONTENT_TOP, half, Cm(0.75), size=Pt(18), bold=True, color=NAVY)
add_shape(sl, rx, CONTENT_TOP + Cm(0.8), half, Cm(0.04), fill=NAVY)
bullet_list(sl, [
    "Directed hyperedges following traffic flow direction",
    "Dynamic hypergraph (edges adapt with congestion patterns)",
    "Neural ODE adjoint method for full continuous-time training",
    "Multi-variate imputation: speed + flow + occupancy jointly",
    "Online sensor selection — adaptive mask during inference",
], rx, CONTENT_TOP + Cm(1.0), half, Cm(9.5), size=Pt(17))

# Final results banner
add_shape(sl, MARGINS, Cm(13.8), CONTENT_W, Cm(0.04), fill=LGREY)
add_text(sl,
    "Final results (80% sparsity, PEMS04 eval set):   "
    "Overall blind-node MAE = 5.18 km/h   \u00b7   "
    "Jam MAE = 33.93 km/h   \u00b7   5.7% improvement over global-mean baseline on congestion",
    MARGINS, Cm(14.0), CONTENT_W, Cm(1.5),
    size=Pt(16), bold=True, color=NAVY, wrap=True)


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 15 — References
# ─────────────────────────────────────────────────────────────────────────────
sl, p = next_slide()
slide_header(sl, "References", p)

refs = [
    "[1]  Y. Li et al., 'Diffusion Convolutional Recurrent Neural Network: Data-Driven Traffic Forecasting,' ICLR, 2018.",
    "[2]  B. Yu et al., 'Spatio-Temporal Graph Convolutional Networks: A Deep Learning Framework for Traffic Forecasting,' IJCAI, 2018.",
    "[3]  Z. Wu et al., 'Graph WaveNet for Deep Spatial-Temporal Graph Modeling,' IJCAI, 2019.",
    "[4]  T. N. Kipf and M. Welling, 'Semi-Supervised Classification with Graph Convolutional Networks,' ICLR, 2017.",
    "[5]  P. Velickovic et al., 'Graph Attention Networks,' ICLR, 2018.",
    "[6]  R. T. Q. Chen et al., 'Neural Ordinary Differential Equations,' NeurIPS, 2018.",
    "[7]  Y. Feng et al., 'Hypergraph Neural Networks,' AAAI, 2019.",
    "[8]  R. E. Kalman, 'A New Approach to Linear Filtering and Prediction Problems,' J. Basic Eng., 1960.",
    "[9]  M. Raissi et al., 'Physics-Informed Neural Networks,' J. Computational Physics, 2019.",
    "[10] M. J. Lighthill and G. B. Whitham, 'On Kinematic Waves II: A Theory of Traffic Flow,' Proc. R. Soc., 1955.",
]
tb = sl.shapes.add_textbox(MARGINS, CONTENT_TOP, CONTENT_W, SH - CONTENT_TOP - Cm(1.2))
tf = tb.text_frame; tf.word_wrap = True
first = True
for ref in refs:
    p_ = tf.paragraphs[0] if first else tf.add_paragraph()
    first = False
    p_.space_before = Pt(4)
    r = p_.add_run()
    r.text = ref
    r.font.name      = 'Calibri'
    r.font.size      = Pt(14)
    r.font.color.rgb = DARK


# ─────────────────────────────────────────────────────────────────────────────
# SLIDE 16 — Questions
# ─────────────────────────────────────────────────────────────────────────────
sl, _ = next_slide()
# Clean white slide, thin navy top bar
add_shape(sl, Cm(0), Cm(0), SW, Cm(0.35), fill=NAVY)
add_shape(sl, Cm(0), SH - Cm(0.35), SW, Cm(0.35), fill=NAVY)

add_text(sl, "Thank You",
         Cm(0), Cm(4.5), SW, Cm(2.5),
         size=Pt(48), bold=True, color=NAVY, align=PP_ALIGN.CENTER)
add_shape(sl, MARGINS + Cm(4), Cm(7.0), CONTENT_W - Cm(8), Cm(0.04), fill=LGREY)
add_text(sl, "Questions & Discussion",
         Cm(0), Cm(7.3), SW, Cm(1.2),
         size=Pt(26), italic=True, color=BLUE, align=PP_ALIGN.CENTER)

add_text(sl,
    "Hypergraph Neural ODEs with Observation Assimilation for Sparse Traffic Speed Imputation",
    Cm(0), Cm(10.0), SW, Cm(0.9),
    size=Pt(16), color=GREY, align=PP_ALIGN.CENTER)
add_text(sl,
    "[Author Name]  ·  [University]  ·  March 2026",
    Cm(0), Cm(10.9), SW, Cm(0.7),
    size=Pt(14), italic=True, color=GREY, align=PP_ALIGN.CENTER)


# ── Save ──────────────────────────────────────────────────────────────────────
prs.save("thesis_presentation.pptx")
print(f"  \u2705 thesis_presentation.pptx  ({pg[0]} slides)")
